#!/usr/bin/env python3
"""Generate Hermes Agent Architecture PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
PRIMARY = RGBColor(0x1a, 0x1a, 0x2e)       # Deep navy
SECONDARY = RGBColor(0x16, 0x21, 0x3e)     # Dark blue
ACCENT = RGBColor(0x0f, 0x34, 0x60)         # Medium blue
HIGHLIGHT = RGBColor(0xe9, 0x45, 0x60)      # Coral red
LIGHT_BG = RGBColor(0xf8, 0xf9, 0xfa)       # Light gray
MEDIUM_GRAY = RGBColor(0x6c, 0x75, 0x7d)
DARK_TEXT = RGBColor(0x21, 0x25, 0x29)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_BLUE = RGBColor(0xe8, 0xf4, 0xf8)
LIGHT_GREEN = RGBColor(0xe8, 0xf5, 0xe9)
LIGHT_ORANGE = RGBColor(0xff, 0xf3, 0xe0)
LIGHT_PURPLE = RGBColor(0xf3, 0xe5, 0xf5)
LIGHT_PINK = RGBColor(0xfc, 0xe4, 0xec)
LIGHT_TEAL = RGBColor(0xe0, 0xf2, 0xf1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(tf, text, font_size=14, bold=False, color=DARK_TEXT, space_after=6, font_name="Microsoft YaHei"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.space_after = Pt(space_after)
    return p

def title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, PRIMARY)
    
    # Accent bar
    add_shape(slide, Inches(0), Inches(2.8), Inches(0.15), Inches(1.8), HIGHLIGHT)
    
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1.2), title, font_size=44, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.8), Inches(10), Inches(0.8), subtitle, font_size=22, color=RGBColor(0xcc, 0xcc, 0xcc))
    
    # Bottom accent
    add_shape(slide, Inches(4), Inches(5.5), Inches(5), Inches(0.05), HIGHLIGHT)
    add_textbox(slide, Inches(4), Inches(5.7), Inches(5), Inches(0.5), "May 2026  |  Xiaomiao AI", font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

def section_slide(section_num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), HIGHLIGHT)
    add_textbox(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.8), section_num, font_size=60, bold=True, color=HIGHLIGHT)
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5), title, font_size=36, bold=True, color=WHITE)
    add_shape(slide, Inches(1), Inches(4.2), Inches(4), Inches(0.05), HIGHLIGHT)

def content_slide(title, bullets, bg_color=LIGHT_BG, accent=ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, bg_color)
    
    # Header bar
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.9), title, font_size=28, bold=True, color=WHITE)
    
    # Content area
    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(11.5)
    height = Inches(5.2)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, (subtitle, items) in enumerate(bullets):
        if subtitle:
            p = tf.paragraphs[0] if i == 0 and not items else tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = accent
            p.font.name = "Microsoft YaHei"
            p.space_after = Pt(8)
        
        if isinstance(items, list):
            for item in items:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Microsoft YaHei"
                p.space_after = Pt(4)
        elif items:
            p = tf.add_paragraph()
            p.text = items
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT
            p.font.name = "Microsoft YaHei"
            p.space_after = Pt(4)
    
    return slide

def two_column_slide(title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.9), title, font_size=28, bold=True, color=WHITE)
    
    # Left column
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.5), WHITE, RGBColor(0xdd, 0xee, 0xe6))
    add_textbox(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.6), left_title, font_size=20, bold=True, color=ACCENT)
    
    txBox_l = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(5.5), Inches(4.5))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    for item in left_items:
        p = tf_l.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(5)
    
    # Right column
    add_shape(slide, Inches(6.9), Inches(1.5), Inches(5.9), Inches(5.5), WHITE, RGBColor(0xdd, 0xee, 0xe6))
    add_textbox(slide, Inches(7.1), Inches(1.6), Inches(5.5), Inches(0.6), right_title, font_size=20, bold=True, color=ACCENT)
    
    txBox_r = slide.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.5), Inches(4.5))
    tf_r = txBox_r.text_frame
    tf_r.word_wrap = True
    for item in right_items:
        p = tf_r.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(5)

def arch_diagram_slide():
    """Architecture diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.9), "Six-Layer Architecture", font_size=28, bold=True, color=WHITE)
    
    layers = [
        ("L6  Infrastructure", "cron/  hermes_cli/  config.yaml", LIGHT_TEAL),
        ("L5  Prompt Building", "agent/prompt_builder.py", LIGHT_PINK),
        ("L4  Tool Execution", "tools/  toolsets.py  registry.py", LIGHT_PURPLE),
        ("L3  Agent Core", "run_agent.py  agent/  (12,002 lines)", LIGHT_ORANGE),
        ("L2  Session Management", "gateway/session.py  hermes_state.py  state.db", LIGHT_GREEN),
        ("L1  Message Intake", "gateway/platforms/  (20+ platforms)", LIGHT_BLUE),
    ]
    
    box_h = Inches(0.85)
    gap = Inches(0.12)
    start_top = Inches(1.5)
    box_w = Inches(11)
    left = Inches(1.2)
    
    for i, (name, detail, color) in enumerate(layers):
        top = start_top + i * (box_h + gap)
        add_shape(slide, left, top, box_w, box_h, color)
        
        # Layer number badge
        add_shape(slide, left, top, Inches(1.5), box_h, ACCENT)
        add_textbox(slide, left, top + Inches(0.1), Inches(1.5), box_h - Inches(0.2), name.split("  ")[0], font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        
        add_textbox(slide, left + Inches(1.7), top + Inches(0.05), Inches(4), box_h - Inches(0.1), name.split("  ")[1] if "  " in name else "", font_size=14, bold=True, color=DARK_TEXT)
        add_textbox(slide, left + Inches(1.7), top + Inches(0.4), Inches(9), Inches(0.4), detail, font_size=11, color=MEDIUM_GRAY)
    
    # Arrow labels on right
    arrows = ["Schedule/Persist", "Assemble Prompt", "Execute Tools", "Reason+Act", "Route Message", "Receive/Send"]
    for i, label in enumerate(arrows):
        top = start_top + i * (box_h + gap) + Inches(0.25)
        add_textbox(slide, left + box_w + Inches(0.2), top, Inches(1.5), Inches(0.4), label, font_size=9, color=ACCENT, alignment=PP_ALIGN.LEFT)

def flow_diagram_slide():
    """Data flow diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.9), "Message Processing Flow", font_size=28, bold=True, color=WHITE)
    
    steps = [
        ("1", "Platform Adapter", "Receives message from platform"),
        ("2", "Session Manager", "Restores conversation context"),
        ("3", "Agent Core", "Builds system prompt + calls LLM"),
        ("4", "Tool Executor", "Dispatches and executes tool calls"),
        ("5", "Context Manager", "Compresses if needed, updates history"),
        ("6", "Response Delivery", "Routes response back to platform"),
    ]
    
    box_w = Inches(1.8)
    box_h = Inches(1.4)
    gap = Inches(0.35)
    start_left = Inches(0.5)
    top = Inches(2)
    
    for i, (num, title, desc) in enumerate(steps):
        left = start_left + i * (box_w + gap)
        
        # Box
        add_shape(slide, left, top, box_w, box_h, WHITE, ACCENT)
        
        # Number circle
        add_shape(slide, left + Inches(0.7), top - Inches(0.25), Inches(0.5), Inches(0.5), HIGHLIGHT)
        add_textbox(slide, left + Inches(0.7), top - Inches(0.2), Inches(0.5), Inches(0.4), num, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        
        # Title
        add_textbox(slide, left + Inches(0.1), top + Inches(0.15), box_w - Inches(0.2), Inches(0.5), title, font_size=11, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
        
        # Description
        add_textbox(slide, left + Inches(0.1), top + Inches(0.6), box_w - Inches(0.2), Inches(0.7), desc, font_size=8, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
        
        # Arrow (not after last)
        if i < len(steps) - 1:
            arrow_left = left + box_w + Inches(0.05)
            add_textbox(slide, arrow_left, top + Inches(0.4), Inches(0.25), Inches(0.5), "->", font_size=20, bold=True, color=HIGHLIGHT, alignment=PP_ALIGN.CENTER)
    
    # Loop arrow at bottom
    add_shape(slide, Inches(1), top + box_h + Inches(0.5), Inches(11), Inches(0.05), HIGHLIGHT)
    add_textbox(slide, Inches(5), top + box_h + Inches(0.7), Inches(3), Inches(0.5), "Loop until LLM stops requesting tools", font_size=13, color=HIGHLIGHT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# GENERATE SLIDES
# ═══════════════════════════════════════════

# Slide 1: Title
title_slide("Hermes Agent", "Architecture Deep Dive & Design Philosophy")

# Slide 2: TOC
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), "Agenda", font_size=32, bold=True, color=WHITE)
items = [
    "Project Overview & Positioning",
    "Six-Layer Architecture Design",
    "Core Module Deep Dive",
    "Agent Runtime Loop",
    "Tool System & MCP Integration",
    "Context Compression Engine",
    "Memory & Session Management",
    "Gateway & Multi-Platform Integration",
    "Error Classification & Failover",
    "Skills System",
    "Design Philosophy Summary",
]
for i, item in enumerate(items):
    col = i // 6
    row = i % 6
    left = Inches(0.8 + col * 6)
    top = Inches(1.3 + row * 0.85)
    add_textbox(slide, left, top, Inches(5.5), Inches(0.6), f"  {i+1}. {item}", font_size=16, color=WHITE)

# Slide 3: Section 1
section_slide("01", "Project Overview")

# Slide 4: What is Hermes
content_slide("What is Hermes Agent?", [
    ("", ["Open-source AI Agent framework by Nous Research", "Named after Greek messenger god Hermes", "Production-grade AI Assistant Runtime", "Supports 20+ messaging platforms", "Core capabilities: tool calling, persistent memory, scheduled tasks, context compression, multi-model switching"]),
])

# Slide 5: Key Capabilities
two_column_slide(
    "Key Capabilities",
    "Core Features",
    [
        "Multi-platform Gateway (20+)",
        "Tool Calling Loop with parallel execution",
        "Persistent Memory with FTS5 search",
        "Smart Context Compression",
        "Cron Job Scheduling",
    ],
    "Advanced Features",
    [
        "Multi-model Support (OpenAI, Anthropic, Google, Ollama)",
        "Smart Error Recovery (12+ error types)",
        "Skills System (20+ categories)",
        "MCP Integration (Model Context Protocol)",
        "Sub-agent Delegation",
    ]
)

# Slide 6: Section 2
section_slide("02", "Architecture Design")

# Slide 7: Architecture Diagram
arch_diagram_slide()

# Slide 8: Data Flow
flow_diagram_slide()

# Slide 9: Design Patterns
content_slide("Design Patterns Used", [
    ("Creational", ["Singleton: ToolRegistry ensures global unique registry", "Factory: Platform adapters via factory methods"]),
    ("Behavioral", ["Strategy: ContextEngine supports multiple strategies", "Observer: asyncio event loop for concurrent messages", "Chain of Responsibility: Error classifier priority chain", "Template Method: run_agent.py conversation loop skeleton"]),
    ("Structural", ["Decorator: Tool registration via registry.register() decorator"]),
])

# Slide 10: Section 3
section_slide("03", "Core Modules")

# Slide 11: run_agent.py
content_slide("run_agent.py -- Agent Main Loop (12,002 lines)", [
    ("Role", "Core engine implementing ReAct (Reasoning + Acting) loop"),
    ("Key Components", [
        "SafeWriter: Wraps stdout/stderr to prevent pipe breakage crashes",
        "jittered_backoff: Exponential backoff + jitter to avoid thundering herd",
        "classify_api_error: Categorizes API errors into 12+ types",
        "Tool Budget: Prevents infinite tool calling loops",
    ]),
])

# Slide 12: model_tools.py & registry.py
two_column_slide(
    "Tool Orchestration & Registry",
    "model_tools.py (611 lines)",
    [
        "get_tool_definitions() -- JSON schemas for LLM",
        "handle_function_call() -- Routes calls to handlers",
        "check_toolset_requirements() -- Validates prerequisites",
        "Thin orchestration layer between Agent and tools",
    ],
    "registry.py (482 lines)",
    [
        "Thread-safe singleton with RLock",
        "ToolEntry: name, toolset, schema, handler, env",
        "discover_builtin_tools() via AST analysis",
        "_snapshot_* methods for stable concurrent reads",
    ]
)

# Slide 13: Section 4
section_slide("04", "Agent Runtime Loop")

# Slide 14: Conversation Loop
content_slide("Conversation Loop Steps", [
    ("The Loop", [
        "1. Receive message from session manager",
        "2. Build system prompt (identity + platform + memory + skills + context)",
        "3. Inject memory via MemoryManager",
        "4. Call LLM with message history + tool definitions",
        "5. Parse response, check finish_reason",
        "6. Execute tool calls in parallel",
        "7. Append results to conversation history",
        "8. Check context compression threshold",
        "9. Loop back or return final response",
    ]),
])

# Slide 15: Error Recovery
content_slide("Error Recovery Mechanism", [
    ("Error Classification", [
        "auth/auth_permanent: Rotate credentials or abort",
        "rate_limit: Backoff then rotate credentials",
        "billing: Rotate immediately",
        "overloaded/server_error: Backoff retry",
        "timeout: Rebuild client then retry",
        "context_overflow: Trigger compression",
        "model_not_found: Fallback to different model",
    ]),
])

# Slide 16: Section 5
section_slide("05", "Tool System")

# Slide 17: Tool System Overview
content_slide("Tool System Overview", [
    ("ToolEntry Structure", ["name, toolset, schema (JSON Schema), handler (async), check_fn, requires_env, is_async, description, emoji"]),
    ("Core Categories", [
        "File Operations: read_file, write_file, patch, search_files",
        "Terminal: terminal, process",
        "Browser: navigate, click, snapshot, type, scroll, vision",
        "Code: execute_code, delegate_task",
        "Memory: memory, session_search",
        "Skills: skill_view, skill_manage, skills_list",
        "Messaging: send_message, text_to_speech",
        "Scheduling: cronjob",
    ]),
])

# Slide 18: MCP Integration
content_slide("MCP (Model Context Protocol) Integration", [
    ("", [
        "Supports dynamic external tool server registration",
        "MCP tools integrate via mcp_tool.py",
        "Registry supports dynamic refresh",
        "When MCP servers add/remove tools, registry auto-updates",
        "New tool definitions exposed to LLM automatically",
    ]),
])

# Slide 19: Section 6
section_slide("06", "Context Compression")

# Slide 20: Compression Algorithm
content_slide("Context Compression Algorithm", [
    ("6-Step Process", [
        "1. Tool Output Pruning (pre-pass): Replace old large outputs with summaries",
        "2. Protect Head: System prompt + first 3 messages kept intact",
        "3. Protect Tail: Recent ~20K tokens protected",
        "4. Deduplication: Merge duplicate tool outputs",
        "5. LLM Summarization: Structured template for middle messages",
        "6. Iterative Update: Update previous summaries, don't regenerate",
    ]),
])

# Slide 21: Anti-Bounce & Template
two_column_slide(
    "Anti-Bounce Protection",
    "Anti-Bounce",
    [
        "Skip compression if last 2 compressions saved < 10%",
        "Suggest /new to start fresh session",
        "600-second cooldown on compression failures",
        "Prevents repeated attempts when service unavailable",
    ],
    "Summary Template",
    [
        "Prefix: CONTEXT COMPACTION -- REFERENCE ONLY",
        "Completed Work: What was done",
        "Current State: Where we are now",
        "Unresolved Issues: Open problems",
        "Remaining Work: What's left to do",
    ]
)

# Slide 22: Section 7
section_slide("07", "Memory & Session")

# Slide 23: Memory System
content_slide("Memory System Architecture", [
    ("Two Layers", [
        "Semantic Memory: memory tool stores user preferences, environmental facts",
        "Episodic Memory: SQLite database stores complete session history with FTS5 search",
    ]),
    ("MemoryManager", [
        "Orchestrates built-in + one external plugin provider",
        "Fence tags: <memory-context> with 'NOT new user input' annotation",
        "Prevents memory pollution from recalled memories",
    ]),
    ("Supported Plugins", ["mem0, Hindsight, Supermemory, ReteroDB, OpenViking, Holographic Memory, Honcho, ByteRover"]),
])

# Slide 24: Session Management
content_slide("Session Management", [
    ("Session Source", ["platform, chat_id, user_id, chat_type (DM/group/channel/thread), thread_id"]),
    ("Key Features", [
        "PII Hashing: SHA-256 hash of user/chat IDs (first 12 hex chars)",
        "WAL Mode: Multi-reader + single-writer concurrency",
        "FTS5: Cross-session full-text search",
        "parent_session_id chain for session continuity",
    ]),
    ("Reset Strategies", ["Manual: /new command", "Idle timeout: configurable", "Context overflow: auto-trigger with summary"]),
])

# Slide 25: Section 8
section_slide("08", "Gateway & Platforms")

# Slide 26: Gateway
content_slide("Gateway Architecture", [
    ("Supported Platforms (20+)", ["Telegram, Discord, WhatsApp, Signal, Slack, Mattermost, Matrix", "WeChat, WeCom, Feishu, DingTalk, QQ Bot, Yuanbao", "iMessage (BlueBubbles), Email, SMS, Home Assistant, API Server, Webhook"]),
    ("Lifecycle", [
        "1. SSL certificate detection (supports NixOS)",
        "2. Load ~/.hermes/.env",
        "3. LRU Agent cache (max 128, 1h TTL)",
        "4. Start platform adapter coroutines",
        "5. Message routing loop",
        "6. Graceful shutdown on SIGTERM",
    ]),
])

# Slide 27: Section 9
section_slide("09", "Error Handling")

# Slide 28: Error Classification
content_slide("Error Classification & Failover", [
    ("12+ Error Types", [
        "auth: 401/403 -> refresh/rotate credentials",
        "rate_limit: 429 -> backoff + rotate",
        "billing: 402 -> rotate immediately",
        "overloaded: 503/529 -> backoff retry",
        "timeout: connection/read timeout -> rebuild client",
        "context_overflow: compress context, not failover",
        "payload_too_large: 413 -> compress payload",
        "model_not_found: 404 -> fallback model",
        "format_error: 400 -> abort or fix",
        "thinking_signature: Anthropic sig invalid",
        "unknown: backoff retry",
    ]),
])

# Slide 29: Section 10
section_slide("10", "Skills System")

# Slide 30: Skills
content_slide("Skills System", [
    ("Discovery", ["iter_skill_index_files() scans skill directories", "Reads YAML frontmatter from SKILL.md", "Injects skill index into system prompt"]),
    ("Loading", ["skill_view() loads full skill content when matched", "LRU cache (max 8 entries)", "Auto-invalidates when skill files change"]),
    ("20+ Categories", ["autonomous-ai-agents, creative, data-science, devops, doc-design", "email, gaming, github, mcp, media, mlops, note-taking", "productivity, red-teaming, research, smart-home, social-media", "software-development, web-access, yuanbao"]),
])

# Slide 31: Section 11
section_slide("11", "Design Philosophy")

# Slide 32: Architecture Highlights
content_slide("Architecture Highlights", [
    ("", [
        "Clear Layering: Six single-responsibility layers",
        "Defensive Programming: SafeWriter, security scanning, PII hashing",
        "Resilient Design: 12+ error types with auto-recovery",
        "Zero-config Extension: AST-based tool auto-discovery",
        "Platform Agnostic: Unified abstraction for 20+ platforms",
        "Persistence First: SQLite WAL + FTS5, all state persisted",
        "Context Aware: Smart compression with fence tags",
        "Observability: Token tracking, cost estimation, detailed logging",
    ]),
])

# Slide 33: Best Practices
content_slide("Best Practices", [
    ("", [
        "All prompt-building functions are stateless (pure functions)",
        "RLock (not Lock) for registry -- supports same-thread reentrancy",
        "WAL mode SQLite + application-level random jitter retry",
        "Context summaries marked 'REFERENCE ONLY'",
        "Tool budget prevents infinite loops, graceful degradation",
        "Skills use lazy loading -- only loaded when needed",
        "Env vars from ~/.hermes/.env override shell exports",
    ]),
])

# Slide 34: Reusable Patterns
content_slide("Reusable Patterns for AI Agent Builders", [
    ("", [
        "Registry Pattern: decorator/self-registration for auto tool discovery",
        "Fence Tags: XML tags to distinguish prompt content types",
        "Structured Error Classification: unified classifier + recovery strategy",
        "Progressive Context Compression: prune first (no LLM), then summarize (with LLM)",
        "Multi-Platform Adapter: unified interface + platform-specific implementation",
    ]),
])

# Slide 35: Ending
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), HIGHLIGHT)
add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2), "Thank You!", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8), "Questions & Discussion", font_size=24, color=RGBColor(0xcc, 0xcc, 0xcc), alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(4.5), Inches(5), Inches(4), Inches(0.05), HIGHLIGHT)
add_textbox(slide, Inches(4), Inches(5.3), Inches(5), Inches(0.5), "by Xiaomiao AI  |  May 2026", font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ── Save ──
out = os.path.expanduser("~/.hermes/Hermes_Architecture.pptx")
prs.save(out)
print(f"PPT generated: {out}")

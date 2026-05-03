#!/usr/bin/env python3
"""生成 Hermes Agent 架构 PPT - 中文版"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

PRIMARY = RGBColor(0x1a, 0x1a, 0x2e)
SECONDARY = RGBColor(0x16, 0x21, 0x3e)
ACCENT = RGBColor(0x0f, 0x34, 0x60)
HIGHLIGHT = RGBColor(0xe9, 0x45, 0x60)
LIGHT_BG = RGBColor(0xf8, 0xf9, 0xfa)
MEDIUM_GRAY = RGBColor(0x6c, 0x75, 0x7d)
DARK_TEXT = RGBColor(0x21, 0x25, 0x29)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_BLUE = RGBColor(0xe8, 0xf4, 0xf8)
LIGHT_GREEN = RGBColor(0xe8, 0xf5, 0xe9)
LIGHT_ORANGE = RGBColor(0xff, 0xf3, 0xe0)
LIGHT_PURPLE = RGBColor(0xf3, 0xe5, 0xf5)
LIGHT_PINK = RGBColor(0xfc, 0xe4, 0xec)
LIGHT_TEAL = RGBColor(0xe0, 0xf2, 0xf1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rounded_box(slide, left, top, width, height, fill_color, text, font_size=14, font_color=WHITE, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].font.color.rgb = font_color
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    return shape

def add_card(slide, left, top, width, height, bg_color, title, body_text, title_color=WHITE, body_color=DARK_TEXT, font_name="Microsoft YaHei"):
    add_shape(slide, left, top, width, height, bg_color)
    add_shape(slide, left, top, width, Pt(3), ACCENT)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                title, font_size=16, bold=True, color=title_color, font_name=font_name)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(0.7),
                body_text, font_size=11, color=body_color, font_name=font_name)

def title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_shape(slide, Inches(0.5), Inches(2.5), Inches(12.333), Inches(0.08), HIGHLIGHT)
    add_textbox(slide, Inches(1), Inches(2.8), Inches(11.333), Inches(1),
                title, font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.9), Inches(11.333), Inches(0.6),
                subtitle, font_size=20, color=RGBColor(0xaa, 0xaa, 0xaa), alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.04), RGBColor(0x44, 0x44, 0x44))
    add_textbox(slide, Inches(1), Inches(5.2), Inches(11.333), Inches(0.5),
                "2026年5月 | 小喵出品", font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

def section_slide(number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11.333), Inches(0.8),
                number, font_size=72, bold=True, color=HIGHLIGHT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.5), Inches(11.333), Inches(1),
                title, font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

def content_slide(title, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
                title, font_size=28, bold=True, color=WHITE)
    for i, bullet in enumerate(bullets):
        y = Inches(1.5) + i * Inches(0.55)
        add_shape(slide, Inches(0.8), y + Inches(0.1), Inches(0.08), Inches(0.08), HIGHLIGHT)
        add_textbox(slide, Inches(1.1), y, Inches(11.5), Inches(0.5),
                    bullet, font_size=16, color=DARK_TEXT)
    if notes:
        add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.733), Inches(0.5),
                    notes, font_size=11, color=MEDIUM_GRAY)

def two_col_slide(title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
                title, font_size=28, bold=True, color=WHITE)

    add_shape(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(5.8), WHITE, RGBColor(0xdd, 0xdd, 0xdd))
    add_shape(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(0.5), ACCENT)
    add_textbox(slide, Inches(0.7), Inches(1.45), Inches(5.6), Inches(0.4),
                left_title, font_size=18, bold=True, color=WHITE)

    add_shape(slide, Inches(6.833), Inches(1.4), Inches(6), Inches(5.8), WHITE, RGBColor(0xdd, 0xdd, 0xdd))
    add_shape(slide, Inches(6.833), Inches(1.4), Inches(6), Inches(0.5), HIGHLIGHT)
    add_textbox(slide, Inches(7.033), Inches(1.45), Inches(5.6), Inches(0.4),
                right_title, font_size=18, bold=True, color=WHITE)

    for i, item in enumerate(left_items):
        y = Inches(2.1) + i * Inches(0.5)
        add_textbox(slide, Inches(0.7), y, Inches(5.6), Inches(0.45),
                    f"• {item}", font_size=13, color=DARK_TEXT)

    for i, item in enumerate(right_items):
        y = Inches(2.1) + i * Inches(0.5)
        add_textbox(slide, Inches(7.033), y, Inches(5.6), Inches(0.45),
                    f"• {item}", font_size=13, color=DARK_TEXT)

def end_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_shape(slide, Inches(0.5), Inches(2.5), Inches(12.333), Inches(0.08), HIGHLIGHT)
    add_textbox(slide, Inches(1), Inches(3), Inches(11.333), Inches(1),
                "谢谢观看", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(11.333), Inches(0.6),
                "Q & A", font_size=24, color=RGBColor(0xaa, 0xaa, 0xaa), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.5),
                "小喵 AI | 2026年5月", font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ═══ 第1页：封面 ═══
title_slide("Hermes Agent", "架构深度解析与设计哲学")

# ═══ 第2页：目录 ═══
section_slide("目录", "内容概览")
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
            "目录", font_size=28, bold=True, color=WHITE)

toc_left = [
    "01 项目概览",
    "02 架构设计",
    "03 核心模块",
    "04 运行时循环",
    "05 工具系统",
    "06 提示词构建",
    "07 上下文压缩",
    "08 记忆系统",
]
toc_right = [
    "09 会话管理",
    "10 网关与平台",
    "11 定时任务",
    "12 错误分类",
    "13 配置系统",
    "14 技能系统",
    "15 设计哲学",
    "",
]

for i, item in enumerate(toc_left):
    y = Inches(1.5) + i * Inches(0.55)
    add_textbox(slide, Inches(0.8), y, Inches(5.5), Inches(0.45),
                item, font_size=16, color=DARK_TEXT)

for i, item in enumerate(toc_right):
    y = Inches(1.5) + i * Inches(0.55)
    add_textbox(slide, Inches(6.833), y, Inches(5.5), Inches(0.45),
                item, font_size=16, color=DARK_TEXT)

# ═══ 第3页：01 项目概览 ═══
section_slide("01", "项目概览")

# ═══ 第4页：核心定位 ═══
content_slide("核心定位",
    [
        "Hermes Agent 是由 Nous Research 开发的开源 AI Agent 框架",
        "以希腊信使神 Hermes 命名，象征信息传递与连接的核心使命",
        "生产级 AI Agent 系统，支持 20+ 消息平台交互",
        "具备工具调用、持久化记忆、定时任务、上下文压缩等能力",
        "被设计为通用 AI 助手运行时，而非简单的聊天机器人框架",
    ],
    "定位：从消息接入到响应投递的全流水线 Agent 操作系统")

# ═══ 第5页：核心能力 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
            "核心能力", font_size=28, bold=True, color=WHITE)

features = [
    ("多平台网关", "20+ 消息平台统一接入", LIGHT_BLUE),
    ("工具调用循环", "ReAct 模式自主执行", LIGHT_GREEN),
    ("持久化记忆", "SQLite + FTS5 全文搜索", LIGHT_ORANGE),
    ("上下文压缩", "智能摘要解决长对话溢出", LIGHT_PURPLE),
    ("定时任务", "Cron 表达式 + 间隔调度", LIGHT_PINK),
    ("多模型支持", "OpenAI / Anthropic / Google 统一接口", LIGHT_TEAL),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + col * Inches(6.333)
    top = Inches(1.5) + row * Inches(1.8)
    add_card(slide, left, top, Inches(6), Inches(1.5), color, title, desc, title_color=DARK_TEXT, body_color=DARK_TEXT)

# ═══ 第6页：技术栈 ═══
content_slide("技术栈",
    [
        "语言：Python 3.12+",
        "核心依赖：OpenAI SDK（兼容多后端）、SQLite、aiohttp、python-dotenv",
        "代码规模：约 12,000 行",
        "模块组织：agent/、gateway/、tools/、cron/、hermes_cli/",
        "架构模式：六层模块化，单职责原则",
    ],
    "代码库地址：github.com/NousResearch/hermes-agent")

# ═══ 第7页：02 架构设计 ═══
section_slide("02", "架构设计")

# ═══ 第8页：六层架构 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
            "六层架构", font_size=28, bold=True, color=WHITE)

layers = [
    ("L1", "消息接入", "多平台消息收发", LIGHT_BLUE),
    ("L2", "会话管理", "会话生命周期、消息持久化", LIGHT_GREEN),
    ("L3", "Agent 核心", "对话循环、工具调用、错误恢复", LIGHT_ORANGE),
    ("L4", "工具执行", "工具定义、注册、执行分发", LIGHT_PURPLE),
    ("L5", "提示词构建", "系统提示词组装、上下文注入", LIGHT_PINK),
    ("L6", "基础设施", "定时任务、CLI、配置管理", LIGHT_TEAL),
]

for i, (num, title, desc, color) in enumerate(layers):
    left = Inches(0.5) + (i % 3) * Inches(4.2)
    top = Inches(1.5) + (i // 3) * Inches(2.7)
    add_shape(slide, left, top, Inches(3.8), Inches(2.3), color)
    add_textbox(slide, left, top + Inches(0.1), Inches(3.8), Inches(0.4),
                num, font_size=14, bold=True, color=HIGHLIGHT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.5), Inches(3.8), Inches(0.5),
                title, font_size=18, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), top + Inches(1.1), Inches(3.4), Inches(1),
                desc, font_size=12, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ═══ 第9页：数据流 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
            "核心数据流", font_size=28, bold=True, color=WHITE)

flow = [
    ("平台适配器", "接收消息"),
    ("会话管理器", "恢复上下文"),
    ("Agent 核心", "构建提示词"),
    ("LLM 推理", "生成工具调用"),
    ("工具执行器", "调度执行"),
    ("结果返回", "追加到历史"),
    ("循环/返回", "直至完成"),
    ("平台适配器", "发送响应"),
]

for i, (name, desc) in enumerate(flow):
    left = Inches(0.3) + i * Inches(1.6)
    top = Inches(2.5)
    add_rounded_box(slide, left, top, Inches(1.4), Inches(1.2), ACCENT, name, font_size=11)
    if i < len(flow) - 1:
        add_textbox(slide, left + Inches(1.45), top + Inches(0.4), Inches(0.15), Inches(0.3),
                    "→", font_size=20, bold=True, color=HIGHLIGHT, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.5),
            "完整流程：接收 → 恢复上下文 → 构建提示词 → LLM 推理 → 工具执行 → 结果追加 → 循环 → 响应",
            font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(5), Inches(12.333), Inches(0.5),
            "每个环节都有完善的错误处理和恢复机制",
            font_size=12, color=HIGHLIGHT, alignment=PP_ALIGN.CENTER)

# ═══ 第10页：设计模式 ═══
two_col_slide("设计模式",
    "创建型 & 结构型",
    [
        "单例模式：ToolRegistry 全局唯一",
        "工厂模式：平台适配器工厂创建",
        "装饰器模式：register() 装饰器注册",
        "模板方法模式：Agent 执行骨架",
    ],
    "行为型",
    [
        "策略模式：多上下文处理策略",
        "观察者模式：asyncio 事件循环",
        "责任链模式：错误分类器优先级链",
    ])

# ═══ 第11页：03 核心模块 ═══
section_slide("03", "核心模块")

# ═══ 第12页：run_agent.py ═══
content_slide("run_agent.py — Agent 主循环",
    [
        "12,002 行，实现完整 ReAct（推理 + 行动）循环",
        "流程：接收消息 → 构建提示词 → 调用 LLM → 解析工具调用 → 执行工具 → 追加结果 → 循环",
        "SafeWriter：包装 stdout/stderr 防止管道断裂崩溃",
        "jittered_backoff：指数退避 + 抖动避免惊群效应",
        "classify_api_error：API 错误分类为 12+ 种类型",
    ],
    "这是整个系统最核心的文件")

# ═══ 第13页：model_tools.py & registry.py ═══
two_col_slide("工具编排层 & 注册中心",
    "model_tools.py（611 行）",
    [
        "get_tool_definitions()：返回所有工具 Schema",
        "handle_function_call()：路由工具调用",
        "check_toolset_requirements()：验证前置条件",
        "轻量编排层，连接 Agent 和工具",
    ],
    "registry.py（482 行）",
    [
        "线程安全的单例工具注册中心",
        "ToolEntry：名称、Schema、处理器、环境需求",
        "threading.RLock 并发保护",
        "AST 分析自动发现内置工具",
    ])

# ═══ 第14页：04 运行时循环 ═══
section_slide("04", "Agent 运行时循环")

# ═══ 第15页：对话循环步骤 ═══
content_slide("对话循环的 9 个步骤",
    [
        "① 接收消息：从会话管理器获取当前消息和对话历史",
        "② 构建提示词：组装系统提示词（身份 + 平台 + 记忆 + 技能）",
        "③ 注入记忆：通过 MemoryManager 预取相关记忆",
        "④ LLM 推理：调用 OpenAI 兼容 API 生成响应",
        "⑤ 解析响应：检查 finish_reason，判断是否需要工具调用",
        "⑥ 执行工具：并行执行所有工具调用，收集结果",
        "⑦ 追加结果：将工具结果添加到对话历史",
        "⑧ 检查压缩：上下文超阈值时触发压缩",
        "⑨ 循环/返回：继续循环或返回最终响应",
    ],
    "循环直到 LLM 不再请求工具调用或达到预算上限")

# ═══ 第16页：错误恢复 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
            "错误恢复策略", font_size=28, bold=True, color=WHITE)

strategies = [
    ("认证失败", "401/403", "刷新/轮换凭证或中止", LIGHT_PINK),
    ("速率限制", "429", "退避后轮换凭证", LIGHT_ORANGE),
    ("服务器错误", "500/503", "退避重试", LIGHT_BLUE),
    ("上下文溢出", "超限", "触发上下文压缩", LIGHT_PURPLE),
    ("超时", "Timeout", "重建客户端后重试", LIGHT_GREEN),
]

for i, (name, status, strategy, color) in enumerate(strategies):
    left = Inches(0.5) + (i % 3) * Inches(4.2)
    top = Inches(1.5) + (i // 3) * Inches(2.5)
    add_card(slide, left, top, Inches(3.8), Inches(2.1), color,
             f"{name} ({status})", strategy, title_color=DARK_TEXT, body_color=DARK_TEXT)

# ═══ 第17页：05 工具系统 ═══
section_slide("05", "工具系统")

# ═══ 第18页：工具分类 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
            "核心工具分类（30+ 内置工具）", font_size=28, bold=True, color=WHITE)

categories = [
    ("文件操作", "read_file, write_file, patch, search_files"),
    ("终端", "terminal, process（后台进程管理）"),
    ("Web 搜索", "web_search, web_extract（网页抓取）"),
    ("浏览器自动化", "navigate, click, snapshot, type..."),
    ("代码执行", "execute_code, delegate_task（子 Agent）"),
    ("视觉", "vision_analyze, image_generate"),
    ("记忆", "memory, session_search"),
    ("任务规划", "todo, clarify"),
    ("消息", "send_message, text_to_speech"),
    ("技能", "skill_view, skill_manage, skills_list"),
    ("调度", "cronjob（定时任务）"),
]

for i, (cat, tools) in enumerate(categories):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + col * Inches(6.333)
    top = Inches(1.4) + row * Inches(0.52)
    colors = [LIGHT_BLUE, LIGHT_GREEN, LIGHT_ORANGE, LIGHT_PURPLE, LIGHT_PINK, LIGHT_TEAL]
    add_shape(slide, left, top, Inches(6), Inches(0.45), colors[i % len(colors)])
    add_textbox(slide, left + Inches(0.2), top + Inches(0.02), Inches(5.6), Inches(0.4),
                f"{cat}: {tools}", font_size=12, color=DARK_TEXT)

# ═══ 第19页：MCP 集成 ═══
content_slide("MCP 集成 — 模型上下文协议",
    [
        "支持 Model Context Protocol（MCP）标准",
        "动态外部工具服务器注册",
        "MCP 工具通过 mcp_tool.py 集成",
        "注册中心支持动态刷新",
        "当 MCP 服务器添加/删除工具时自动更新",
        "新工具定义自动暴露给 LLM",
    ],
    "MCP 让 Hermes 能够接入任何兼容标准的外部工具生态")

# ═══ 第20页：06 提示词构建 ═══
section_slide("06", "提示词构建系统")

# ═══ 第21页：系统提示词组件 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
            "系统提示词组件", font_size=28, bold=True, color=WHITE)

components = [
    ("身份定义", "DEFAULT_AGENT_IDENTITY — 基础人格和行为准则"),
    ("平台提示", "PLATFORM_HINTS — 17+ 平台特定交互指南"),
    ("记忆引导", "MEMORY_GUIDANCE — 记忆系统使用指南"),
    ("会话搜索", "SESSION_SEARCH_GUIDANCE — 历史会话搜索指南"),
    ("技能引导", "SKILLS_GUIDANCE — 技能使用和创建指南"),
    ("工具执行强化", "TOOL_USE_ENFORCEMENT_GUIDANCE — 强制使用工具"),
    ("模型特定", "执行指导因模型系列而异"),
    ("上下文文件", "AGENTS.md、SOUL.md、.hermes.md"),
    ("环境提示", "WSL / Docker 运行时环境感知"),
    ("语言检测", "自动检测用户语言并调整响应语言"),
]

for i, (name, desc) in enumerate(components):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + col * Inches(6.333)
    top = Inches(1.4) + row * Inches(0.52)
    add_textbox(slide, left, top, Inches(6), Inches(0.45),
                f"• {name}: {desc}", font_size=12, color=DARK_TEXT)

# ═══ 第22页：安全扫描 ═══
content_slide("上下文文件安全扫描",
    [
        "在注入上下文文件前执行安全扫描",
        "检测不可见 Unicode 字符（零宽字符、双向覆盖）",
        "检测威胁模式（如「忽略之前的指令」）",
        "威胁内容被拦截并替换为警告",
        "防止提示词注入攻击",
    ],
    "安全扫描是防止恶意上下文文件的关键防线")

# ═══ 第23页：07 上下文压缩 ═══
section_slide("07", "上下文压缩系统")

# ═══ 第24页：压缩流程 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
            "压缩算法流程", font_size=28, bold=True, color=WHITE)

steps = [
    ("① 工具输出剪枝", "用摘要替换旧的大型工具输出（无需 LLM）"),
    ("② 保护头部", "系统提示词 + 前 N 条消息保持完整"),
    ("③ 保护尾部", "按 token 预算保护最近消息（~20K）"),
    ("④ 去重", "合并重复的工具输出"),
    ("⑤ LLM 摘要", "结构化模板摘要中间消息"),
    ("⑥ 迭代更新", "更新之前的摘要而非重新生成"),
]

for i, (title, desc) in enumerate(steps):
    left = Inches(0.5) + (i % 2) * Inches(6.333)
    top = Inches(1.4) + (i // 2) * Inches(1.6)
    add_card(slide, left, top, Inches(6), Inches(1.3), LIGHT_BLUE, title, desc, title_color=DARK_TEXT, body_color=DARK_TEXT)

# ═══ 第25页：防回弹 ═══
content_slide("防回弹保护",
    [
        "问题：压缩循环每次仅压缩 1-2 条消息，效率极低",
        "方案：最近两次压缩节省率均 < 10% 时跳过压缩",
        "建议用户 /new 开始新会话",
        "600 秒冷却时间防止摘要服务不可用时重复尝试",
        "结构化摘要模板：已完成工作、当前状态、未解决问题、剩余工作",
        '摘要前缀："CONTEXT COMPACTION -- REFERENCE ONLY"',
    ],
    "防回弹确保压缩真正有效，而非无效循环")

# ═══ 第26页：08 记忆系统 ═══
section_slide("08", "记忆系统")

# ═══ 第27页：双层记忆 ═══
two_col_slide("双层记忆架构",
    "语义记忆",
    [
        "通过 memory 工具存储",
        "用户偏好和环境事实",
        "跨会话持久化",
        "SQLite + FTS5 全文搜索",
        "memory_manager.py（373 行）协调",
    ],
    "情节记忆",
    [
        "通过 hermes_state.py 存储",
        "完整会话历史记录",
        "FTS5 虚拟表全文搜索",
        "parent_session_id 链支持会话分裂",
        "WAL 模式支持并发读写",
    ])

# ═══ 第28页：记忆围栏 ═══
content_slide("记忆围栏设计",
    [
        "记忆上下文通过 <memory-context> 围栏标签注入",
        '附带系统注释："NOT new user input"',
        "防止模型将召回的记忆当作新用户输入",
        "有效防止记忆污染",
        "支持插件：mem0、Hindsight、Supermemory 等 8 种",
    ],
    "围栏设计是记忆系统的关键安全机制")

# ═══ 第29页：09 会话管理 ═══
section_slide("09", "会话管理")

# ═══ 第30页：会话管理 ═══
content_slide("会话管理（session.py — 1,257 行）",
    [
        "SessionSource：平台、chat_id、user_id、chat_type",
        "支持类型：私信（DM）、群组、频道、话题（thread）",
        "每条消息携带完整追踪信息用于路由",
        "PII 哈希：user_id/chat_id 经 SHA-256 哈希后存储",
        "多种重置策略：手动 /new、空闲超时、上下文溢出",
        "parent_session_id 链在重置间保持连续性",
    ],
    "会话管理是多平台消息路由的核心")

# ═══ 第31页：10 网关与平台 ═══
section_slide("10", "网关与平台集成")

# ═══ 第32页：网关 ═══
content_slide("网关架构（gateway/run.py — 11,015 行）",
    [
        "管理所有平台适配器的生命周期",
        "asyncio 事件循环处理多平台并发消息",
        "每个平台适配器运行在各自的协程中",
        "支持 20+ 平台：Telegram、Discord、微信、飞书等",
        "统一适配器接口：connect、listen、send、disconnect",
        "SSL 证书自动检测（支持 NixOS）",
        "LRU 缓存最多 128 个 Agent 实例",
    ],
    "网关是多平台接入的核心基础设施")

# ═══ 第33页：11 定时任务 ═══
section_slide("11", "定时任务系统")

# ═══ 第34页：定时任务 ═══
content_slide("定时任务系统（cron/jobs.py — 776 行）",
    [
        "任务存储：~/.hermes/cron/jobs.json",
        "输出保存：~/.hermes/cron/output/{job_id}/",
        "三种调度类型：一次性、间隔、Cron 表达式",
        "多种投递模式：origin、local、指定平台",
        "一次性任务 120 秒宽限期",
        "间隔任务宽限期为间隔的一半（120s ~ 2h）",
        "threading.Lock 保护并发读写",
    ],
    "定时任务让 Agent 具备了自主定时执行能力")

# ═══ 第35页：12 错误分类 ═══
section_slide("12", "错误分类与故障转移")

# ═══ 第36页：错误分类 ═══
content_slide("错误分类体系（error_classifier.py — 829 行）",
    [
        "结构化 API 错误分类，替代分散的字符串匹配",
        "12+ 种错误类型：认证、速率限制、计费、过载、超时等",
        "每种类型有针对性恢复策略",
        "ClassifiedError 封装：reason、status_code、provider、model",
        "恢复提示：retryable、should_compress、should_rotate_credential",
        "主重试循环直接检查恢复提示",
    ],
    "集中式错误分类器是整个恢复系统的决策中枢")

# ═══ 第37页：13 配置系统 ═══
section_slide("13", "配置系统")

# ═══ 第38页：配置 ═══
content_slide("多层配置系统",
    [
        "优先级：CLI 参数 > 环境变量 > config.yaml",
        "config.yaml 主配置：~/.hermes/config.yaml",
        "章节：model、gateway、tools、session、cron、display、plugins",
        "环境变量前缀：HERMES_（如 HERMES_MODEL=gpt-4）",
        "支持多模型提供者：OpenAI、Anthropic、Google、Ollama 等",
        "故障转移时自动切换到备用模型",
    ],
    "多层配置确保灵活性和可维护性的平衡")

# ═══ 第39页：14 技能系统 ═══
section_slide("14", "技能系统")

# ═══ 第40页：技能系统 ═══
content_slide("技能系统",
    [
        "可插拔的能力扩展机制",
        "技能 = SKILL.md 文件（YAML frontmatter + Markdown 正文）",
        "存储位置：~/.hermes/skills/ 目录",
        "自动发现：扫描目录下所有 SKILL.md 文件",
        "9 种类型、5 种模式",
        "技能缓存基于文件路径和修改时间",
    ],
    "技能系统让 Agent 具备了无限扩展的能力")

# ═══ 第41页：15 设计哲学 ═══
section_slide("15", "设计哲学总结")

# ═══ 第42页：架构亮点 ═══
content_slide("架构亮点",
    [
        "六层模块化：每层职责单一，接口清晰",
        "ReAct 循环：Reasoning + Acting 自主分解复杂任务",
        "智能错误恢复：12+ 种错误分类，针对性恢复策略",
        "上下文压缩：保护头尾 + LLM 摘要中间历史",
        "双层记忆：语义记忆 + 情节记忆，FTS5 跨会话搜索",
        "多平台统一：20+ 平台共享统一核心",
        "声明式工具注册：零配置集成新工具",
    ],
    "简洁而不简单，灵活而不混乱，强大而不臃肿")

# ═══ 第43页：最佳实践 ═══
content_slide("最佳实践",
    [
        "优先使用已有工具，而非编写新代码",
        "长时间任务必须使用任务队列并定期汇报进度",
        "错误恢复时先分类再处理，避免盲目重试",
        "上下文接近阈值时主动触发压缩，而非等待溢出",
        "记忆存储用户偏好和环境事实，不存储临时任务状态",
    ],
    "这些实践是从实际使用中总结的经验")

# ═══ 第44页：结尾 ═══
end_slide()

prs.save('/home/ubuntu/.hermes/Hermes_Architecture_CN.pptx')
print("PPT 生成成功: /home/ubuntu/.hermes/Hermes_Architecture_CN.pptx")
